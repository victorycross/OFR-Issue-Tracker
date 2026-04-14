#!/usr/bin/env python3
"""
Create the OFR Issue Deck template PPTX from a corporate source template.

Opens the source PPTX, deletes all content slides (retaining the slide master
and all 69 layouts), and saves a blank template for use by generate_issue_deck.py.

This is a one-time setup script. The output template retains the PwC Office
theme (Georgia/Arial fonts, PwC colour scheme) and all slide layouts.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Emu

BASE = os.path.dirname(os.path.abspath(__file__))
DEFAULT_SOURCE = os.path.expanduser(
    "~/Downloads/AI Governance/Deployment Patterns and Transition Requirements v1 Draft.pptx"
)
DEFAULT_OUTPUT = os.path.join(BASE, "templates", "OFR_Issue_Deck_Template.pptx")

# Layout indices we use in generate_issue_deck.py (for reference)
LAYOUT_MAP = {
    0:  "Cover (Title Slide)",
    4:  "Issue Table (Title and Content)",
    9:  "Section Divider (Section Header)",
    18: "Executive Summary (Agenda 2)",
    22: "Issue Detail / Split (Two Content)",
    26: "Group Summary (Three Content)",
    30: "KPI Dashboard (Data and Four Stats)",
    58: "Status Overview (Four Stat Boxes)",
    64: "Flexible / Appendix (Title Only)",
    66: "Closing (Conclusion 1)",
}


def main():
    source = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_SOURCE
    output = sys.argv[2] if len(sys.argv) > 2 else DEFAULT_OUTPUT

    if not os.path.exists(source):
        print(f"Error: Source template not found: {source}")
        sys.exit(1)

    prs = Presentation(source)

    # Report template info
    w_in = prs.slide_width / Emu(914400)
    h_in = prs.slide_height / Emu(914400)
    print(f"Source: {source}")
    print(f"Slide dimensions: {w_in:.2f}\" x {h_in:.2f}\"")
    print(f"Slide masters: {len(prs.slide_masters)}")
    print(f"Layouts in master: {len(prs.slide_masters[0].slide_layouts)}")
    print(f"Content slides: {len(prs.slides)}")

    # Delete all content slides using the XML API for clean removal
    slide_count = len(prs.slides)
    sldIdLst = prs.slides._sldIdLst
    sldId_elements = list(sldIdLst)  # snapshot to avoid mutation during iteration

    for sldId in sldId_elements:
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        # Remove the relationship and the slide part
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    print(f"\nDeleted {slide_count} content slides")

    # Verify clean state
    remaining = len(list(prs.slides._sldIdLst))
    print(f"Remaining slides: {remaining}")

    # Print layout map for reference
    print(f"\nLayouts retained (all {len(prs.slide_masters[0].slide_layouts)}):")
    print(f"Layouts we use ({len(LAYOUT_MAP)}):")
    for idx, purpose in sorted(LAYOUT_MAP.items()):
        layout = prs.slide_masters[0].slide_layouts[idx]
        print(f"  [{idx:2d}] {layout.name:<30s} -> {purpose}")

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output), exist_ok=True)

    prs.save(output)
    print(f"\nTemplate saved: {output}")
    print(f"File size: {os.path.getsize(output):,} bytes")


if __name__ == "__main__":
    main()
