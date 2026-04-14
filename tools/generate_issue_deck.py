#!/usr/bin/env python3
"""
Generate the OFR Issue Tracker — Weekly Risk Report PPTX.

Uses the OFR template (cleaned from the corporate PwC slide master) to produce
a presentation summarizing all open issues, grouped by functional area, with
KPI dashboards, staleness indicators, and colour-coded priority/status tables.

Usage:
    python3 generate_issue_deck.py -i data.json [-o output.pptx] [-t template.pptx]

The input JSON schema is documented in sample-data/ofr_issues_sample.json.
"""

import argparse
import json
import os
import sys
from collections import Counter
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = os.path.dirname(os.path.abspath(__file__))
DEFAULT_TEMPLATE = os.path.join(BASE, "templates", "OFR_Issue_Deck_Template.pptx")

# ── Layout indices (from slide master) ──────────────────────────────────────
LY_COVER            = 0   # Title Slide — CENTER_TITLE + SUBTITLE
LY_ISSUE_TABLE      = 4   # Title and Content — TITLE + OBJECT
LY_SECTION_DIVIDER  = 9   # Section Header — TITLE + BODY
LY_EXEC_SUMMARY     = 18  # Agenda 2 — TITLE + BODY
LY_SPLIT            = 22  # Two Content — TITLE + 2x OBJECT
LY_THREE_COL        = 26  # Three Content — TITLE + 3x OBJECT
LY_KPI_DASHBOARD    = 30  # Data and Four Stats — TITLE + BODY + 4x OBJECT
LY_TITLE_ONLY       = 64  # Title Only — TITLE
LY_CLOSING          = 66  # Conclusion 1 — CENTER_TITLE + SUBTITLE

# ── Appkit4 Design System Colours ───────────────────────────────────────────
APPKIT_BLUE   = RGBColor(0x41, 0x53, 0x85)  # Primary Blue — headers, Current, Low, New
APPKIT_ORANGE = RGBColor(0xD0, 0x4A, 0x02)  # Primary Orange — Active
APPKIT_RED    = RGBColor(0xE0, 0x30, 0x1E)  # Primary Red — High, Escalated, Stale
APPKIT_AMBER  = RGBColor(0xE4, 0x5C, 0x2B)  # Orange +0.5 — Medium, Monitoring, Aging
NEUTRAL_BLACK = RGBColor(0x2D, 0x2D, 0x2D)  # Body text
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_600      = RGBColor(0x64, 0x74, 0x8B)
GRAY_400      = RGBColor(0x94, 0xA3, 0xB8)
GRAY_LIGHT    = RGBColor(0xF5, 0xF5, 0xF5)  # Alternate row shading
GRAY_BORDER   = RGBColor(0xE1, 0xE1, 0xE1)  # Table borders

# ── Functional Group Display Order ──────────────────────────────────────────
GROUP_ORDER = [
    "Risk Management Office",
    "Engagement Risk",
    "Client Risk and KYC",
    "Technology Risk & AI Trust",
    "OGC Contracts",
    "OGC Privacy",
    "Internal Audit",
    "National Security",
    "OGC General Counsel",
    "Independence",
]

# Short names for compact display
GROUP_SHORT = {
    "Risk Management Office": "RMO",
    "Engagement Risk": "EngRisk",
    "Client Risk and KYC": "ClientRisk",
    "Technology Risk & AI Trust": "TechRisk",
    "OGC Contracts": "OGCContract",
    "OGC Privacy": "OGCPrivacy",
    "Internal Audit": "IntAudit",
    "National Security": "NatSec",
    "OGC General Counsel": "OGCCounsel",
    "Independence": "Indep",
}


# ══════════════════════════════════════════════════════════════════════════════
#  COLOUR HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def color_for_priority(priority):
    return {"High": APPKIT_RED, "Medium": APPKIT_AMBER, "Low": APPKIT_BLUE}.get(priority, GRAY_400)

def color_for_status(status):
    return {
        "New": APPKIT_BLUE, "Active": APPKIT_ORANGE,
        "Escalated": APPKIT_RED, "Monitoring": APPKIT_AMBER,
        "Closed": GRAY_400,
    }.get(status, GRAY_400)

def color_for_staleness(flag):
    return {"Current": APPKIT_BLUE, "Aging": APPKIT_AMBER, "Stale": APPKIT_RED}.get(flag, GRAY_400)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def add_slide(prs, layout_idx):
    """Add a slide using the specified layout index."""
    layout = prs.slide_masters[0].slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def get_ph(slide, idx):
    """Safely get a placeholder by index, or None if not found.

    NOTE: SlidePlaceholders does not support 'in' for index lookup — the 'in'
    operator iterates over placeholder objects, not indices. Use try/except instead.
    """
    try:
        return slide.placeholders[idx]
    except KeyError:
        return None


def set_text(placeholder, text, font_size=None, bold=False, color=None, alignment=None):
    """Set text on a placeholder's first paragraph with formatting."""
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    if font_size:
        run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color
    return run


def add_paragraph(text_frame, text, font_size=11, bold=False, color=None, alignment=None, space_before=0, space_after=0):
    """Add a new paragraph to a text frame."""
    p = text_frame.add_paragraph()
    if alignment:
        p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color
    return p


def set_cell_text(cell, text, font_size=9, bold=False, color=None, alignment=None):
    """Set text in a table cell with formatting."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color


def set_cell_fill(cell, color):
    """Set the fill colour of a table cell."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def add_table_at(slide, rows, cols, left, top, width, height):
    """Add a table to the slide at the specified position."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


# ══════════════════════════════════════════════════════════════════════════════
#  DATA HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def compute_stats(issues):
    """Compute summary statistics from the issue list."""
    open_issues = [i for i in issues if i["Status"] != "Closed"]
    stats = {
        "total_open": len(open_issues),
        "total_closed": len(issues) - len(open_issues),
        "high": sum(1 for i in open_issues if i["Priority"] == "High"),
        "medium": sum(1 for i in open_issues if i["Priority"] == "Medium"),
        "low": sum(1 for i in open_issues if i["Priority"] == "Low"),
        "new": sum(1 for i in open_issues if i["Status"] == "New"),
        "active": sum(1 for i in open_issues if i["Status"] == "Active"),
        "escalated": sum(1 for i in open_issues if i["Status"] == "Escalated"),
        "monitoring": sum(1 for i in open_issues if i["Status"] == "Monitoring"),
        "current": sum(1 for i in open_issues if i["StalenessFlag"] == "Current"),
        "aging": sum(1 for i in open_issues if i["StalenessFlag"] == "Aging"),
        "stale": sum(1 for i in open_issues if i["StalenessFlag"] == "Stale"),
        "avg_days": round(sum(i["DaysSinceUpdate"] for i in open_issues) / max(len(open_issues), 1), 1),
        "groups_with_issues": len(set(i["FunctionalGroup"] for i in open_issues)),
    }
    return stats


def issues_by_group(issues):
    """Group open issues by functional group, in display order."""
    open_issues = [i for i in issues if i["Status"] != "Closed"]
    grouped = {}
    for group in GROUP_ORDER:
        group_issues = [i for i in open_issues if i["FunctionalGroup"] == group]
        if group_issues:
            grouped[group] = sorted(group_issues, key=lambda x: -x["DaysSinceUpdate"])
    return grouped


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def build_cover(prs, data):
    """Slide 1: Cover page."""
    slide = add_slide(prs, LY_COVER)
    title_text = data.get("report_title", "One Firm Risk Tracker")
    subtitle_text = data.get("report_subtitle", f"Weekly Issue Report")

    # Title Slide has CENTER_TITLE at idx 0 and SUBTITLE at idx 1
    ph = get_ph(slide, 0)
    if ph:
        set_text(ph, title_text, font_size=32, bold=True)
    ph = get_ph(slide, 1)
    if ph:
        set_text(ph, subtitle_text, font_size=16, color=GRAY_600)


def build_executive_summary(prs, data, stats):
    """Slide 2: Executive Summary narrative."""
    slide = add_slide(prs, LY_EXEC_SUMMARY)

    ph = get_ph(slide, 0)
    if ph:
        set_text(ph, "Executive Summary", font_size=28, bold=True)

    # Build the body narrative
    body_ph = None
    for idx in [12, 1, 13]:  # Agenda 2 body placeholder varies
        body_ph = get_ph(slide, idx)
        if body_ph:
            break

    if body_ph:
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True

        # Summary paragraph
        summary = (
            f"As of {data.get('generated_at', 'today')[:10]}, the OFR Issue Tracker has "
            f"{stats['total_open']} open items across {stats['groups_with_issues']} functional groups."
        )
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = summary
        run.font.size = Pt(12)
        run.font.color.rgb = NEUTRAL_BLACK

        # Key findings
        findings = [
            f"High Priority: {stats['high']} items require immediate attention",
            f"Stale Items: {stats['stale']} items have not been updated in 15+ days",
            f"Aging Items: {stats['aging']} items are between 8-14 days since last update",
            f"Average days since update: {stats['avg_days']}",
        ]
        for finding in findings:
            add_paragraph(tf, f"\u2022  {finding}", font_size=11, color=NEUTRAL_BLACK, space_before=4)

        # Escalations
        if stats['escalated'] > 0:
            add_paragraph(tf, "", font_size=6)  # spacer
            add_paragraph(
                tf,
                f"\u26a0  {stats['escalated']} item(s) currently escalated",
                font_size=12, bold=True, color=APPKIT_RED, space_before=8
            )


def build_kpi_dashboard(prs, stats):
    """Slide 3: KPI Dashboard with four stat boxes."""
    slide = add_slide(prs, LY_TITLE_ONLY)

    ph = get_ph(slide, 0)
    if ph:
        set_text(ph, "Issue Portfolio at a Glance", font_size=28, bold=True)

    # Since the Data and Four Stats layout placeholder indices can be unreliable,
    # we use Title Only and add shapes manually for maximum control.
    kpis = [
        ("Open Items", str(stats["total_open"]), APPKIT_BLUE),
        ("Stale (15+ days)", str(stats["stale"]), APPKIT_RED),
        ("High Priority", str(stats["high"]), APPKIT_RED if stats["high"] > 0 else APPKIT_BLUE),
        ("Avg Days Since Update", str(stats["avg_days"]), APPKIT_AMBER if stats["avg_days"] > 10 else APPKIT_BLUE),
    ]

    box_width = Inches(2.0)
    box_height = Inches(1.5)
    start_x = Inches(0.5)
    start_y = Inches(2.2)
    gap = Inches(0.25)

    for i, (label, value, color) in enumerate(kpis):
        x = start_x + i * (box_width + gap)
        # Add a rounded rectangle
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()  # no border

        tf = shape.text_frame
        tf.word_wrap = True

        # Value (large number)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(16)
        run = p.add_run()
        run.text = value
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Label
        add_paragraph(tf, label, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Second row: status distribution
    statuses = [
        ("New", stats["new"], APPKIT_BLUE),
        ("Active", stats["active"], APPKIT_ORANGE),
        ("Escalated", stats["escalated"], APPKIT_RED),
        ("Monitoring", stats["monitoring"], APPKIT_AMBER),
    ]

    box2_height = Inches(1.0)
    row2_y = start_y + box_height + Inches(0.4)

    for i, (label, value, color) in enumerate(statuses):
        x = start_x + i * (box_width + gap)
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, row2_y, box_width, box2_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = str(value)
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = WHITE

        add_paragraph(tf, label, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


def build_staleness_overview(prs, data, stats):
    """Slide 4: Staleness overview — left breakdown, right stale items list."""
    slide = add_slide(prs, LY_TITLE_ONLY)

    ph = get_ph(slide, 0)
    if ph:
        set_text(ph, "Staleness Overview", font_size=28, bold=True)

    # Left side: staleness breakdown boxes
    categories = [
        ("Current (0-7 days)", stats["current"], APPKIT_BLUE),
        ("Aging (8-14 days)", stats["aging"], APPKIT_AMBER),
        ("Stale (15+ days)", stats["stale"], APPKIT_RED),
    ]

    from pptx.enum.shapes import MSO_SHAPE
    left_x = Inches(0.5)
    y_start = Inches(2.0)
    box_w = Inches(4.0)
    box_h = Inches(1.0)
    gap = Inches(0.2)

    for i, (label, count, color) in enumerate(categories):
        y = y_start + i * (box_h + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(16)
        run = p.add_run()
        run.text = f"  {count}"
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run2 = p.add_run()
        run2.text = f"   {label}"
        run2.font.size = Pt(14)
        run2.font.color.rgb = WHITE

    # Right side: stale items list
    open_issues = [i for i in data["issues"] if i["Status"] != "Closed"]
    stale_items = sorted(
        [i for i in open_issues if i["StalenessFlag"] == "Stale"],
        key=lambda x: -x["DaysSinceUpdate"]
    )

    if stale_items:
        right_x = Inches(5.0)
        # Add a text box for the stale items list
        txBox = slide.shapes.add_textbox(right_x, y_start, Inches(4.5), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Items Requiring Attention"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = APPKIT_RED

        for item in stale_items[:8]:  # Cap at 8 to fit
            text = f"{item['ItemID']} - {item['Title'][:40]}  ({item['DaysSinceUpdate']}d)"
            add_paragraph(tf, text, font_size=9, color=NEUTRAL_BLACK, space_before=4)
            add_paragraph(
                tf,
                f"   Owner: {item['Owner']} | {item['FunctionalGroup']}",
                font_size=8, color=GRAY_600, space_before=0
            )


def build_group_section(prs, group_name, group_issues):
    """Build a section divider + issue table for one functional group."""
    # Section divider slide
    slide_div = add_slide(prs, LY_SECTION_DIVIDER)
    ph = get_ph(slide_div, 0)
    if ph:
        set_text(ph, group_name, font_size=28, bold=True)

    stale_count = sum(1 for i in group_issues if i["StalenessFlag"] == "Stale")
    summary_text = f"{len(group_issues)} open issue{'s' if len(group_issues) != 1 else ''}"
    if stale_count > 0:
        summary_text += f"  |  {stale_count} stale"

    ph = get_ph(slide_div, 1)
    if ph:
        set_text(ph, summary_text, font_size=14, color=GRAY_600)

    # Issue table slide(s) — paginate at 7 issues per slide
    PAGE_SIZE = 7
    for page_start in range(0, len(group_issues), PAGE_SIZE):
        page_issues = group_issues[page_start:page_start + PAGE_SIZE]
        _build_issue_table_slide(prs, group_name, page_issues, page_start, len(group_issues))


def _build_issue_table_slide(prs, group_name, issues, offset, total):
    """Build a single issue table slide for a group."""
    slide = add_slide(prs, LY_TITLE_ONLY)

    page_label = f" (continued)" if offset > 0 else ""
    short = GROUP_SHORT.get(group_name, group_name)
    ph = get_ph(slide, 0)
    if ph:
        set_text(ph, f"{group_name}{page_label}", font_size=22, bold=True)

    # Table columns: ID, Title, Owner, Priority, Status, Days, Staleness
    columns = ["ID", "Title", "Owner", "Priority", "Status", "Days", "Flag"]
    col_widths = [Inches(0.7), Inches(3.0), Inches(1.4), Inches(0.9), Inches(1.0), Inches(0.6), Inches(0.8)]
    num_rows = len(issues) + 1  # +1 for header

    table_left = Inches(0.5)
    table_top = Inches(1.8)
    table_width = sum(w for w in col_widths)
    table_height = Inches(0.35 * num_rows)

    table = add_table_at(slide, num_rows, len(columns), table_left, table_top, table_width, table_height)

    # Set column widths
    for i, width in enumerate(col_widths):
        table.columns[i].width = width

    # Header row
    for ci, header in enumerate(columns):
        cell = table.cell(0, ci)
        set_cell_text(cell, header, font_size=9, bold=True, color=WHITE)
        set_cell_fill(cell, APPKIT_BLUE)

    # Data rows
    for ri, issue in enumerate(issues, start=1):
        row_data = [
            issue["ItemID"],
            issue["Title"][:45] + ("..." if len(issue["Title"]) > 45 else ""),
            issue["Owner"],
            issue["Priority"],
            issue["Status"],
            str(issue["DaysSinceUpdate"]),
            issue["StalenessFlag"],
        ]

        for ci, value in enumerate(row_data):
            cell = table.cell(ri, ci)
            color = NEUTRAL_BLACK

            # Colour-code specific columns
            if ci == 3:  # Priority
                color = color_for_priority(value)
            elif ci == 4:  # Status
                color = color_for_status(value)
            elif ci == 6:  # Staleness flag
                color = color_for_staleness(value)

            set_cell_text(cell, value, font_size=8, bold=(ci == 0), color=color)

            # Alternate row shading
            if ri % 2 == 0:
                set_cell_fill(cell, GRAY_LIGHT)


def build_action_required(prs, data, stats):
    """Build the 'Items Requiring Attention' slide with all stale items."""
    open_issues = [i for i in data["issues"] if i["Status"] != "Closed"]
    stale_items = sorted(
        [i for i in open_issues if i["StalenessFlag"] == "Stale"],
        key=lambda x: -x["DaysSinceUpdate"]
    )

    if not stale_items:
        return  # Skip if no stale items

    slide = add_slide(prs, LY_TITLE_ONLY)
    ph = get_ph(slide, 0)
    if ph:
        set_text(
            ph,
            f"Items Requiring Attention ({len(stale_items)} Stale)",
            font_size=22, bold=True, color=APPKIT_RED
        )

    columns = ["ID", "Title", "Owner", "Group", "Days", "Next Action"]
    col_widths = [Inches(0.7), Inches(2.2), Inches(1.2), Inches(1.3), Inches(0.6), Inches(2.8)]
    display_items = stale_items[:10]  # Cap at 10
    num_rows = len(display_items) + 1

    table = add_table_at(
        slide, num_rows, len(columns),
        Inches(0.3), Inches(1.8),
        sum(col_widths), Inches(0.35 * num_rows)
    )

    for i, width in enumerate(col_widths):
        table.columns[i].width = width

    # Header
    for ci, header in enumerate(columns):
        cell = table.cell(0, ci)
        set_cell_text(cell, header, font_size=9, bold=True, color=WHITE)
        set_cell_fill(cell, APPKIT_RED)

    # Data
    for ri, issue in enumerate(display_items, start=1):
        short_group = GROUP_SHORT.get(issue["FunctionalGroup"], issue["FunctionalGroup"][:12])
        row_data = [
            issue["ItemID"],
            issue["Title"][:35] + ("..." if len(issue["Title"]) > 35 else ""),
            issue["Owner"],
            short_group,
            str(issue["DaysSinceUpdate"]),
            (issue.get("NextAction") or "")[:40],
        ]
        for ci, value in enumerate(row_data):
            cell = table.cell(ri, ci)
            set_cell_text(cell, value, font_size=8, color=NEUTRAL_BLACK)
            if ri % 2 == 0:
                set_cell_fill(cell, GRAY_LIGHT)


def build_priority_summary(prs, data, priority, header_color):
    """Build a summary table slide for a single priority level (High/Medium/Low)."""
    open_issues = [i for i in data["issues"] if i["Status"] != "Closed"]
    priority_issues = sorted(
        [i for i in open_issues if i["Priority"] == priority],
        key=lambda x: -x["DaysSinceUpdate"]
    )

    if not priority_issues:
        return  # Skip empty priority levels

    # Section divider for the priority band
    slide_div = add_slide(prs, LY_SECTION_DIVIDER)
    ph = get_ph(slide_div, 0)
    if ph:
        set_text(ph, f"{priority} Priority Issues", font_size=28, bold=True, color=header_color)
    ph = get_ph(slide_div, 1)
    if ph:
        stale = sum(1 for i in priority_issues if i["StalenessFlag"] == "Stale")
        esc = sum(1 for i in priority_issues if i["Status"] == "Escalated")
        subtitle_parts = [f"{len(priority_issues)} item{'s' if len(priority_issues) != 1 else ''}"]
        if stale:
            subtitle_parts.append(f"{stale} stale")
        if esc:
            subtitle_parts.append(f"{esc} escalated")
        set_text(ph, "  |  ".join(subtitle_parts), font_size=14, color=GRAY_600)

    # Table slides — paginate at 8 rows per slide
    PAGE_SIZE = 8
    columns = ["ID", "Title", "Owner", "Group", "Status", "Days", "Flag", "Next Action"]
    col_widths = [
        Inches(0.6), Inches(2.0), Inches(1.1), Inches(1.1),
        Inches(0.8), Inches(0.5), Inches(0.6), Inches(2.1),
    ]

    for page_start in range(0, len(priority_issues), PAGE_SIZE):
        page_issues = priority_issues[page_start:page_start + PAGE_SIZE]
        slide = add_slide(prs, LY_TITLE_ONLY)

        page_label = " (continued)" if page_start > 0 else ""
        ph = get_ph(slide, 0)
        if ph:
            set_text(ph, f"{priority} Priority Summary{page_label}", font_size=22, bold=True, color=header_color)

        num_rows = len(page_issues) + 1
        table = add_table_at(
            slide, num_rows, len(columns),
            Inches(0.3), Inches(1.8),
            sum(col_widths), Inches(0.35 * num_rows),
        )

        for ci, width in enumerate(col_widths):
            table.columns[ci].width = width

        # Header row
        for ci, header in enumerate(columns):
            cell = table.cell(0, ci)
            set_cell_text(cell, header, font_size=9, bold=True, color=WHITE)
            set_cell_fill(cell, header_color)

        # Data rows
        for ri, issue in enumerate(page_issues, start=1):
            short_group = GROUP_SHORT.get(issue["FunctionalGroup"], issue["FunctionalGroup"][:12])
            row_data = [
                issue["ItemID"],
                issue["Title"][:30] + ("..." if len(issue["Title"]) > 30 else ""),
                issue["Owner"],
                short_group,
                issue["Status"],
                str(issue["DaysSinceUpdate"]),
                issue["StalenessFlag"],
                (issue.get("NextAction") or "")[:30],
            ]
            for ci, value in enumerate(row_data):
                cell = table.cell(ri, ci)
                color = NEUTRAL_BLACK
                if ci == 4:  # Status
                    color = color_for_status(value)
                elif ci == 6:  # Staleness
                    color = color_for_staleness(value)
                set_cell_text(cell, value, font_size=8, bold=(ci == 0), color=color)
                if ri % 2 == 0:
                    set_cell_fill(cell, GRAY_LIGHT)


def build_issue_detail(prs, issue):
    """Build a formatted single-issue detail slide.

    Layout:
      Title:  "OFR-NNN — Issue Title"
      Row 1:  Priority badge | Status badge | Group | Owner       (colour-coded pills)
      Row 2:  Next Action narrative
      Row 3:  Update History timeline (date + notes, most recent first)
      Footer: Copilot placeholder bar
    """
    from pptx.enum.shapes import MSO_SHAPE

    slide = add_slide(prs, LY_TITLE_ONLY)

    # ── Slide title ──────────────────────────────────────────────────────────
    ph = get_ph(slide, 0)
    if ph:
        set_text(ph, f"{issue['ItemID']}  —  {issue['Title']}", font_size=20, bold=True)

    # ── Badge row: Priority | Status | Group | Owner ─────────────────────────
    badge_y = Inches(1.55)
    badge_h = Inches(0.32)
    gap = Inches(0.12)
    x_cursor = Inches(0.4)

    priority = issue.get("Priority", "Medium")
    status = issue.get("Status", "Active")
    group = issue.get("FunctionalGroup", "")
    owner = issue.get("Owner", "")
    staleness = issue.get("StalenessFlag", "Current")

    badges = [
        (priority, color_for_priority(priority), Inches(1.1)),
        (status, color_for_status(status), Inches(1.1)),
        (staleness, color_for_staleness(staleness), Inches(1.0)),
    ]

    for label, bg_color, width in badges:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_cursor, badge_y, width, badge_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = WHITE
        x_cursor += width + gap

    # Group and Owner as formatted text next to badges
    meta_box = slide.shapes.add_textbox(x_cursor + Inches(0.1), badge_y, Inches(5.5), badge_h)
    tf = meta_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_before = Pt(2)
    # Group
    run = p.add_run()
    run.text = group
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = APPKIT_BLUE
    # Separator
    run = p.add_run()
    run.text = "    \u00b7    "
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY_400
    # Owner
    run = p.add_run()
    run.text = owner
    run.font.size = Pt(10)
    run.font.color.rgb = NEUTRAL_BLACK

    # ── Next Action ──────────────────────────────────────────────────────────
    na_y = badge_y + badge_h + Inches(0.25)
    na_box = slide.shapes.add_textbox(Inches(0.4), na_y, Inches(9.0), Inches(0.7))
    tf = na_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Next Action:  "
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = APPKIT_BLUE
    run = p.add_run()
    run.text = issue.get("NextAction") or "No next action defined"
    run.font.size = Pt(11)
    run.font.color.rgb = NEUTRAL_BLACK

    # Date context line
    raised = issue.get("DateRaised", "")
    updated = issue.get("LastUpdated", "")
    days = issue.get("DaysSinceUpdate", 0)
    add_paragraph(
        tf,
        f"Raised {raised}    \u00b7    Last updated {updated}    \u00b7    {days} days ago",
        font_size=9, color=GRAY_600, space_before=4,
    )

    # ── Update History ────────────────────────────────────────────────────────
    updates = issue.get("Updates") or []
    history_y = na_y + Inches(0.9)

    if updates:
        hist_box = slide.shapes.add_textbox(Inches(0.4), history_y, Inches(9.0), Inches(3.2))
        tf = hist_box.text_frame
        tf.word_wrap = True

        # Section header
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Update History"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = APPKIT_BLUE

        # Show up to 5 most recent updates (already sorted newest first in data)
        for update in updates[:5]:
            u_date = update.get("Date", "")
            u_status = update.get("Status", "")
            u_notes = update.get("Notes", "")
            u_by = update.get("UpdatedBy", "")

            # Date + status line
            p = tf.add_paragraph()
            p.space_before = Pt(8)
            run = p.add_run()
            run.text = f"{u_date}"
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = NEUTRAL_BLACK
            if u_status:
                run = p.add_run()
                run.text = f"  [{u_status}]"
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = color_for_status(u_status)
            if u_by:
                run = p.add_run()
                run.text = f"  — {u_by}"
                run.font.size = Pt(8)
                run.font.color.rgb = GRAY_600

            # Notes line
            if u_notes:
                p = tf.add_paragraph()
                p.space_before = Pt(2)
                run = p.add_run()
                run.text = u_notes
                run.font.size = Pt(9)
                run.font.color.rgb = NEUTRAL_BLACK
    else:
        # No updates — show placeholder
        hist_box = slide.shapes.add_textbox(Inches(0.4), history_y, Inches(9.0), Inches(0.5))
        tf = hist_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "No update history available"
        run.font.size = Pt(10)
        run.font.italic = True
        run.font.color.rgb = GRAY_400

    # ── Copilot executive summary placeholder (footer bar) ───────────────────
    copilot_y = Inches(6.6)
    cp_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), copilot_y, Inches(9.2), Inches(0.6)
    )
    cp_box.fill.solid()
    cp_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFA)  # light blue-grey
    cp_box.line.color.rgb = RGBColor(0xB0, 0xBE, 0xD0)
    cp_box.line.width = Pt(1)

    tf = cp_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = "  \U0001f916  Copilot Executive Brief  "
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = APPKIT_BLUE
    run = p.add_run()
    run.text = (
        "This section will be populated by Microsoft Copilot to provide an "
        "executive-ready summary of this issue for leadership review."
    )
    run.font.size = Pt(8)
    run.font.color.rgb = GRAY_600


def build_closing(prs, data):
    """Final slide: Closing."""
    slide = add_slide(prs, LY_CLOSING)

    ph = get_ph(slide, 0)
    if ph:
        set_text(ph, "One Firm Risk Tracker", font_size=32, bold=True)

    generated = data.get("generated_at", datetime.now().isoformat())[:16].replace("T", " ")
    ph = get_ph(slide, 1)
    if ph:
        set_text(
            ph,
            f"Generated {generated}  \u00b7  OFR Risk & Quality",
            font_size=14, color=GRAY_600
        )


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(description="Generate OFR Issue Deck (PPTX)")
    parser.add_argument("--input", "-i", required=True, help="Path to JSON data file")
    parser.add_argument("--output", "-o", default=None, help="Output PPTX path")
    parser.add_argument("--template", "-t", default=None, help="Template PPTX path")
    args = parser.parse_args()

    # Load data
    with open(args.input) as f:
        data = json.load(f)

    issues = data.get("issues", [])
    if not issues:
        print("Warning: No issues found in input data. Generating minimal deck.")

    # Paths
    template_path = args.template or DEFAULT_TEMPLATE
    if not os.path.exists(template_path):
        print(f"Error: Template not found: {template_path}")
        sys.exit(1)

    output_path = args.output or os.path.join(
        BASE, f"OFR_Issue_Deck_{datetime.now().strftime('%Y-%m-%d')}.pptx"
    )

    # Compute stats
    stats = compute_stats(issues)

    print(f"Input: {args.input} ({len(issues)} issues)")
    print(f"Template: {template_path}")
    print(f"Open: {stats['total_open']} | Stale: {stats['stale']} | High: {stats['high']}")

    # Generate presentation
    prs = Presentation(template_path)

    # Build slides in sequence
    build_cover(prs, data)
    build_executive_summary(prs, data, stats)
    build_kpi_dashboard(prs, stats)
    build_staleness_overview(prs, data, stats)

    # Priority summary tables (High → Medium → Low)
    build_priority_summary(prs, data, "High", APPKIT_RED)
    build_priority_summary(prs, data, "Medium", APPKIT_AMBER)
    build_priority_summary(prs, data, "Low", APPKIT_BLUE)

    # Group sections
    grouped = issues_by_group(issues)
    for group_name, group_issues in grouped.items():
        build_group_section(prs, group_name, group_issues)

    # Individual issue detail slides (one per active issue, sorted by priority)
    open_issues = [i for i in issues if i["Status"] != "Closed"]
    priority_rank = {"High": 0, "Medium": 1, "Low": 2}
    sorted_issues = sorted(open_issues, key=lambda x: (priority_rank.get(x["Priority"], 3), -x["DaysSinceUpdate"]))

    if sorted_issues:
        # Section divider before detail slides
        detail_div = add_slide(prs, LY_SECTION_DIVIDER)
        ph = get_ph(detail_div, 0)
        if ph:
            set_text(ph, "Issue Detail Slides", font_size=28, bold=True)
        ph = get_ph(detail_div, 1)
        if ph:
            set_text(ph, f"{len(sorted_issues)} active issues  ·  one slide per issue", font_size=14, color=GRAY_600)

        for issue in sorted_issues:
            build_issue_detail(prs, issue)

    # Action required (stale items)
    build_action_required(prs, data, stats)

    # Closing
    build_closing(prs, data)

    # Save
    prs.save(output_path)
    slide_count = len(prs.slides)
    file_size = os.path.getsize(output_path)
    print(f"\nIssue Deck PPTX created: {output_path}")
    print(f"Slides: {slide_count}")
    print(f"File size: {file_size:,} bytes")


if __name__ == "__main__":
    main()
